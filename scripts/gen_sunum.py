# -*- coding: utf-8 -*-
"""Finans Portalı — proje sunumu (.pptx) üretici. python-pptx ile."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY   = RGBColor(0x1E, 0x3A, 0x8A)
BLUE   = RGBColor(0x1D, 0x4E, 0xD8)
LBLUE  = RGBColor(0xDB, 0xEA, 0xFE)
INK    = RGBColor(0x1F, 0x29, 0x37)
MUTED  = RGBColor(0x47, 0x55, 0x69)
GREEN  = RGBColor(0x15, 0x80, 0x3D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PALE   = RGBColor(0xF1, 0xF5, 0xF9)

HEAD = "Cambria"
BODY = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, rot=0, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    if rot:
        shp.rotation = rot
    return shp


def accents(s):
    # sağ-üst geometrik köşe (başlık sol-üstte; çakışma olmaz)
    rect(s, Inches(12.5), Inches(-0.8), Inches(2.1), Inches(1.5), NAVY, rot=35)
    rect(s, Inches(12.9), Inches(-0.5), Inches(1.8), Inches(1.3), BLUE, rot=35)
    # sağ-alt köşe
    rect(s, Inches(12.2), Inches(6.6), Inches(2.6), Inches(2.0), NAVY, rot=35)
    rect(s, Inches(12.6), Inches(7.0), Inches(2.4), Inches(1.8), BLUE, rot=35)


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf


def para(tf, text, size=16, color=INK, bold=False, font=BODY, space=6, first=False, align=PP_ALIGN.LEFT, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space); p.level = level
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font; r.font.color.rgb = color
    return p


def title(s, text, sub=None):
    rect(s, Inches(0.85), Inches(0.6), Inches(0.16), Inches(0.95), BLUE)  # accent bar
    tf = tb(s, Inches(1.15), Inches(0.55), Inches(11.2), Inches(1.1))
    para(tf, text, size=34, color=NAVY, bold=True, font=HEAD, first=True, space=2)
    if sub:
        para(tf, sub, size=14, color=MUTED, font=BODY)
    rect(s, Inches(1.15), Inches(1.62), Inches(3.2), Inches(0.05), LBLUE)


# ───────────────────────── 1 · KAPAK ─────────────────────────
s = slide()
rect(s, 0, 0, SW, SH, WHITE)
# büyük navy bant solda
rect(s, Inches(-1.0), Inches(-1.0), Inches(3.0), Inches(9.5), NAVY, rot=12)
rect(s, Inches(-1.2), Inches(-1.0), Inches(2.2), Inches(9.5), BLUE, rot=12)
rect(s, Inches(11.7), Inches(5.9), Inches(3.2), Inches(2.6), NAVY, rot=30)
tf = tb(s, Inches(2.4), Inches(2.2), Inches(9.6), Inches(2.8))
para(tf, "Finans Portalı", size=58, color=NAVY, bold=True, font=HEAD, first=True, space=6)
para(tf, "Çok-varlıklı finans portföy & piyasa takip platformu", size=20, color=MUTED, font=BODY, space=4)
rect(s, Inches(2.45), Inches(4.15), Inches(3.6), Inches(0.06), BLUE)
tf2 = tb(s, Inches(2.4), Inches(4.45), Inches(9.6), Inches(1.8))
para(tf2, "Proje Sunumu", size=22, color=BLUE, bold=True, font=HEAD, first=True, space=10)
para(tf2, "Yiğit Şeker", size=18, color=INK, bold=True, font=BODY, space=2)
para(tf2, "Sakarya Üniversitesi · Yazılım Mühendisliği", size=14, color=MUTED, font=BODY)


# ───────────────────────── 2 · HAKKIMDA ─────────────────────────
s = slide(); accents(s); title(s, "Hakkımda")
tf = tb(s, Inches(1.15), Inches(2.0), Inches(7.2), Inches(4.6))
para(tf, "Ben Yiğit Şeker. Sakarya Üniversitesi Yazılım Mühendisliği 3. sınıf "
         "öğrencisiyim.", size=17, color=INK, first=True, space=12)
para(tf, "Yazılım geliştirme alanında özellikle full-stack web teknolojileri ve "
         "yapay zeka tabanlı bilgisayarlı görü projeleri üzerine çalışıyorum. "
         "Backend, frontend ve veri odaklı sistemler geliştirmenin yanında, "
         "YOLO tabanlı nesne tespiti projeleriyle yapay zekanın gerçek dünya "
         "problemlerine uygulanması konusunda deneyim kazanıyorum.", size=17, color=INK, space=10)
# sağ bilgi kartı
card = rect(s, Inches(8.7), Inches(2.0), Inches(3.6), Inches(3.4), PALE)
ctf = tb(s, Inches(9.0), Inches(2.25), Inches(3.0), Inches(3.0))
para(ctf, "Künye", size=13, color=BLUE, bold=True, font=HEAD, first=True, space=10)
for k, v in [("Bölüm", "Yazılım Müh. · 3. sınıf"), ("Not Ortalaması (GNO)", "2.78"),
             ("İngilizce", "B2"), ("İlgi alanları", "Full-stack web · Yapay zeka / bilgisayarlı görü (YOLO)")]:
    para(ctf, k, size=11, color=MUTED, font=BODY, space=0)
    para(ctf, v, size=14, color=INK, bold=True, font=BODY, space=8)


# ───────────────────────── 3 · PROJE HAKKINDA & AMAÇ ─────────────────────────
s = slide(); accents(s); title(s, "Proje Hakkında & Amaç")
tf = tb(s, Inches(1.15), Inches(1.95), Inches(11.0), Inches(1.6))
para(tf, "Ne sunar?", size=16, color=BLUE, bold=True, font=HEAD, first=True, space=4)
para(tf, "Hisse, kripto, yatırım fonu, tahvil/bono, döviz, emtia ve VİOP'u; enflasyon, faiz ve "
         "döviz kuru gibi makro verileri; finans haberlerini ve kullanıcının kişisel portföyünü "
         "TEK bir modern web arayüzünde birleştirir.", size=15.5, color=INK, space=8)
tf2 = tb(s, Inches(1.15), Inches(3.9), Inches(11.0), Inches(2.6))
para(tf2, "Amaç", size=16, color=BLUE, bold=True, font=HEAD, first=True, space=4)
para(tf2, "Dağınık finans uygulamalarını tek ekranda toplamak ve veriyi yalnızca göstermekle "
          "kalmayıp anlamlandırmak:", size=15.5, color=INK, space=6)
for t in ["Enflasyondan arındırılmış gerçek (reel) getiri — “kazandım mı, gerçekten kazandım mı?”",
          "Risk ve kısa/uzun vade sinyalleri ile çapraz-varlık analiz",
          "Sorumlu yapay zeka danışman (yatırım tavsiyesi değildir uyarısıyla)"]:
    para(tf2, "▸  " + t, size=15, color=INK, space=5)


# ───────────────────────── 4 · ÖNE ÇIKAN ÖZELLİKLER ─────────────────────────
s = slide(); accents(s); title(s, "Öne Çıkan Özellikler")
feats = [
    ("Çok varlıklı canlı piyasa", "Hisse/kripto/fon/tahvil/döviz/emtia/VİOP; grafikler, sparkline, arama-filtre."),
    ("Portföy & reel getiri", "Kâr/zarar, dağılım grafiği, enflasyona göre reel getiri, Excel ile içe aktarma."),
    ("Fiyat alarmları & listeler", "Hedef fiyat alarmı (uygulama içi + e-posta), kişisel takip listeleri."),
    ("Yapay zeka danışman", "Sohbet tabanlı finans asistanı + çapraz-varlık analiz tablosu."),
    ("Tahvil & VİOP simülasyonu", "Long/short, teminat, nominal/kupon/itfa — gerçek emir değil, simülasyon."),
    ("Çok dilli & güvenli", "TR/EN arayüz + otomatik haber çevirisi; 2FA ile güvenli giriş."),
]
x0, y0, cw, ch, gx, gy = Inches(1.15), Inches(2.0), Inches(5.55), Inches(1.5), Inches(0.45), Inches(0.35)
for i, (h, d) in enumerate(feats):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    rect(s, cx, cy, cw, ch, PALE)
    rect(s, cx, cy, Inches(0.09), ch, BLUE)
    ctf = tb(s, cx + Inches(0.28), cy + Inches(0.14), cw - Inches(0.5), ch - Inches(0.2))
    para(ctf, h, size=15, color=NAVY, bold=True, font=HEAD, first=True, space=3)
    para(ctf, d, size=11.5, color=MUTED, font=BODY)


# ───────────────────────── 5 · MİMARİ GENEL BAKIŞ ─────────────────────────
s = slide(); accents(s); title(s, "Mimari Genel Bakış")
# akış kutuları
flow = [("Tarayıcı", "React SPA", LBLUE), ("nginx", "ters vekil", PALE),
        ("Spring Boot", "REST /api/v1", LBLUE), ("PostgreSQL", "Flyway + JPA", PALE)]
fx, fy, fw, fh = Inches(1.15), Inches(2.05), Inches(2.5), Inches(1.0)
for i, (h, d, c) in enumerate(flow):
    cx = fx + (fw + Inches(0.42)) * i
    rect(s, cx, fy, fw, fh, c, line=RGBColor(0xCB,0xD5,0xE1))
    t = tb(s, cx + Inches(0.18), fy + Inches(0.12), fw - Inches(0.3), fh - Inches(0.2))
    para(t, h, size=14, color=NAVY, bold=True, font=HEAD, first=True, space=2)
    para(t, d, size=11, color=MUTED, font=BODY)
    if i < 3:
        a = tb(s, cx + fw, fy, Inches(0.42), fh, anchor=MSO_ANCHOR.MIDDLE)
        para(a, "→", size=20, color=BLUE, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = tb(s, Inches(1.15), Inches(3.5), Inches(11.1), Inches(3.4))
para(tf, "Katmanlı mimari:  Controller → Service → Repository (JPA) → PostgreSQL; "
         "dış kaynaklara erişim ayrı client paketinde izole.", size=14.5, color=INK, bold=True, first=True, space=10)
para(tf, "Kimlik:", size=14, color=BLUE, bold=True, font=HEAD, space=2)
para(tf, "Keycloak (OIDC + 2FA TOTP), Spring Security OAuth2 Resource Server (JWT/RS256).", size=13.5, color=INK, space=8)
para(tf, "Gözlemlenebilirlik:", size=14, color=BLUE, bold=True, font=HEAD, space=2)
para(tf, "Loglar: Log4j2 → Kafka → log-consumer → OpenSearch.   "
         "Metrik: Micrometer → Prometheus → Grafana.   İz: OpenTelemetry → Jaeger.", size=13.5, color=INK, space=8)
para(tf, "Dağıtım:", size=14, color=BLUE, bold=True, font=HEAD, space=2)
para(tf, "Docker Compose (tüm yığın) + Kubernetes / Kustomize / GKE; Caffeine in-memory cache.", size=13.5, color=INK)


# ───────────────────────── 6 · KARŞILANAN İSTERLER ─────────────────────────
s = slide(); accents(s); title(s, "Karşılanan Proje İsterleri", "Zorunlu isterlerin tamamı uygulandı")
groups_l = [
    ("Frontend", "ReactJS"),
    ("Backend", "Java 21 · Spring Boot 3.x · Log4j2"),
    ("Veri & ORM", "PostgreSQL · Spring Data JPA/Hibernate · Flyway"),
    ("Güvenlik", "JWT + Keycloak · 2FA (TOTP)"),
    ("Gözlemlenebilirlik", "OpenTelemetry · Grafana+Prometheus · OpenSearch"),
]
groups_r = [
    ("Log aktarımı", "Kafka (log4j → Kafka → tüketici → OpenSearch)"),
    ("DevOps", "Docker · Git (düzenli commit trafiği)"),
    ("Performans", "Cache — Caffeine (in-memory)"),
    ("Dokümantasyon & Kalite", "REST /api/v1 versioning · OpenAPI/Swagger · Javadocs · README"),
    ("", "Unit Test · Error Handling · Katmanlı Mimari / OOP"),
]
for col, groups in ((Inches(1.15), groups_l), (Inches(7.0), groups_r)):
    t = tb(s, col, Inches(2.05), Inches(5.5), Inches(4.6))
    first = True
    for h, d in groups:
        if h:
            para(t, "✓  " + h, size=15, color=GREEN, bold=True, font=HEAD, first=first, space=1)
        para(t, ("     " if h else "     ") + d, size=12.5, color=INK, space=9, first=(first and not h))
        first = False


# ───────────────────────── 7 · EKSTRA TEKNOLOJİLER ─────────────────────────
s = slide(); accents(s); title(s, "Ekstra Teknolojiler", "İsterlerin ötesinde eklediklerimiz")
extra_l = [
    ("Kubernetes · Kustomize · GKE", "Docker'ın ötesinde üretim orkestrasyonu (WIF ile anahtarsız deploy)"),
    ("GitHub Actions CI/CD", "Otomatik build, test ve dağıtım hattı"),
    ("SonarCloud + JaCoCo", "Kod kalite kapıları + test kapsamı ölçümü"),
    ("Jaeger", "Dağıtık iz (trace) görüntüleme — OTel ile"),
    ("ShedLock", "Dağıtık zamanlayıcı kilidi (çok-replika güvenliği)"),
    ("LLM / Yapay zeka danışman", "OpenAI-uyumlu (Groq) sohbet asistanı"),
]
extra_r = [
    ("LibreTranslate", "Self-hosted TR↔EN haber çevirisi"),
    ("Playwright (headless Chromium)", "Bot-korumalı sayfalardan veri çekimi"),
    ("OpenLDAP", "Keycloak kullanıcı federasyonu"),
    ("Spring Mail + Mailpit", "E-posta fiyat alarmı bildirimleri"),
    ("Apache POI + SheetJS", "Excel ile portföy içe/dışa aktarma"),
    ("Yahoo · TCMB EVDS3 · TEFAS · FRED · İş Yatırım", "Canlı dış veri entegrasyonları"),
]
for col, items in ((Inches(1.15), extra_l), (Inches(7.0), extra_r)):
    t = tb(s, col, Inches(2.05), Inches(5.55), Inches(4.8))
    first = True
    for h, d in items:
        para(t, "▸  " + h, size=14, color=NAVY, bold=True, font=HEAD, first=first, space=1)
        para(t, "     " + d, size=11.5, color=MUTED, font=BODY, space=9)
        first = False


# ───────────────────────── 8 · ZORLUKLAR VE ÖĞRENDİKLERİM ─────────────────────────
s = slide(); accents(s); title(s, "Zorluklar ve Öğrendiklerim")
cards = [
    ("01", "CI/CD Pipeline",
     "Daha önce daha basit projelerimde uyguladığım CI ve CD Pipeline, bu kapsamlı "
     "projede özellikle CD kısmı beni en çok zorlayan bölümlerden oldu."),
    ("02", "Finansal Okuryazarlık",
     "Sistemi yalnızca geliştirici değil, gerçek bir kullanıcının ihtiyaçları üzerinden "
     "değerlendirmeye çalıştım. Başlangıçta finansal kavramları, yatırım araçlarını ve "
     "kullanıcı beklentilerini öğrenmeye odaklandım."),
    ("03", "Yeni Teknolojiler",
     "SonarCloud, JaCoCo ve Grafana gibi araçlar başta zorladı; ancak bu süreç kod kalitesi, "
     "test kapsamı ve sistem izleme konularında bakış açımı geliştirdi."),
]
cw3 = Inches(3.7)
for i, (n, h, d) in enumerate(cards):
    cx = Inches(1.15) + (cw3 + Inches(0.35)) * i
    t = tb(s, cx, Inches(2.1), cw3, Inches(4.6))
    para(t, n, size=40, color=LBLUE, bold=True, font=HEAD, first=True, space=2)
    para(t, h, size=17, color=NAVY, bold=True, font=HEAD, space=8)
    para(t, d, size=12.5, color=INK, space=4)


# ───────────────────────── 9 · EKSİK KALAN İSTERLER ─────────────────────────
s = slide(); accents(s); title(s, "Eksik Kalan İsterler")
big = rect(s, Inches(1.15), Inches(2.05), Inches(11.0), Inches(1.0), RGBColor(0xF0,0xFD,0xF4), line=RGBColor(0xBB,0xF7,0xD0))
t = tb(s, Inches(1.45), Inches(2.2), Inches(10.5), Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
para(t, "✓  Zorunlu isterlerin tamamı (1–24) karşılandı.", size=17, color=GREEN, bold=True, font=HEAD, first=True)
tf = tb(s, Inches(1.15), Inches(3.5), Inches(11.0), Inches(3.0))
para(tf, "Yalnızca opsiyonel madde:", size=15, color=BLUE, bold=True, font=HEAD, first=True, space=6)
para(tf, "▸  Madde 17 — Redis (opsiyonel) kullanılmadı. Yerine Caffeine in-memory cache tercih "
         "edildi; tek-instance için yeterli, çok-instance dağıtımda Redis'e geçiş için altyapı hazır.",
     size=14.5, color=INK, space=10)
para(tf, "Not:", size=15, color=BLUE, bold=True, font=HEAD, space=4)
para(tf, "▸  GKE'ye otomatik CD şu an maliyet nedeniyle manuel tetiklemeye alınmıştır; deploy "
         "hattı çalışır durumdadır.", size=14.5, color=INK)


# ───────────────────────── 10 · KAPANIŞ ─────────────────────────
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, Inches(-1.0), Inches(-1.0), Inches(3.0), Inches(9.5), BLUE, rot=12)
tf = tb(s, Inches(2.4), Inches(2.6), Inches(9.6), Inches(2.6))
para(tf, "Teşekkürler", size=46, color=WHITE, bold=True, font=HEAD, first=True, space=10)
para(tf, "Finans Portalı — çok-varlıklı finans portföy & piyasa takip platformu",
     size=17, color=LBLUE, font=BODY, space=6)
para(tf, "Yiğit Şeker · Sakarya Üniversitesi Yazılım Mühendisliği", size=15, color=WHITE, font=BODY, space=2)
para(tf, "GitHub: github.com/yigitsheker/finans-portali", size=14, color=LBLUE, font=BODY)

prs.save("docs/Finans_Portali_Sunum.pptx")
print("OK", len(prs.slides.__iter__.__self__._sldIdLst), "slayt")
