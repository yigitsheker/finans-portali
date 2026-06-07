# -*- coding: utf-8 -*-
"""Kapanış slaytındaki (10) taşan eğik mavi bandı, köşeye oturan üçgen kamalarla
değiştirir (lacivert zemin korunur, taşma yok). Metne dokunulmaz."""
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PATH  = "docs/Finans_Portali_Sunum.pptx"
BLUE  = RGBColor(0x1D, 0x4E, 0xD8)
LIGHT = RGBColor(0x60, 0xA5, 0xFA)

prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height


def wedge(s, corner, w, h, color):
    W, Hh = Inches(w), Inches(h)
    left = 0 if corner[1] == "L" else SW - W
    top  = 0 if corner[0] == "T" else SH - Hh
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, W, Hh)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    xfrm = shp._element.spPr.find(qn("a:xfrm"))
    if corner[0] == "T": xfrm.set("flipV", "1")
    if corner[1] == "R": xfrm.set("flipH", "1")
    return shp


s = prs.slides[9]  # kapanış
# taşan eğik bandı kaldır (lacivert zemin rot=0 kalır)
for sh in list(s.shapes):
    if abs(getattr(sh, "rotation", 0) or 0) > 0.5:
        sh._element.getparent().remove(sh._element)
# köşeye oturan kamalar (köşelerde; merkezdeki metni örtmez)
wedge(s, "TL", 3.3, 2.3, BLUE); wedge(s, "TL", 2.15, 1.5, LIGHT)
wedge(s, "BR", 3.3, 2.3, BLUE); wedge(s, "BR", 2.15, 1.5, LIGHT)

prs.save(PATH)
print("OK — kapanis kamalari eklendi")
