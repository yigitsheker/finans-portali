# -*- coding: utf-8 -*-
"""Mevcut sunumu YERINDE yamalar: köşe süslerini slayt içine sığdırır (taşma yok)
+ slayt 4/5 kutu alt-yazılarını okunur yapar. Slayt 7 (extra tech) ELLENMEZ."""
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

PATH = "docs/Finans_Portali_Sunum.pptx"
NAVY = RGBColor(0x1E, 0x3A, 0x8A)
BLUE = RGBColor(0x1D, 0x4E, 0xD8)
INK  = RGBColor(0x1F, 0x29, 0x37)

prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height


def add_rect(s, x, y, w, h, color, rot):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False; shp.rotation = rot
    return shp


def contained_accents(s):
    # sağ-üst (slayt içinde kalır)
    add_rect(s, Inches(11.85), Inches(0.45), Inches(1.4), Inches(0.5), NAVY, 35)
    add_rect(s, Inches(12.15), Inches(0.30), Inches(1.1), Inches(0.4), BLUE, 35)
    # sağ-alt
    add_rect(s, Inches(11.85), Inches(6.55), Inches(1.4), Inches(0.5), NAVY, 35)
    add_rect(s, Inches(12.15), Inches(6.78), Inches(1.1), Inches(0.4), BLUE, 35)


# 1) İçerik slaytlarındaki (2–9) eski taşan köşe süslerini kaldır, sığan yenisini ekle
for idx in range(1, 9):
    s = prs.slides[idx]
    for sh in list(s.shapes):
        if abs(getattr(sh, "rotation", 0) or 0) > 0.5:   # döndürülmüş = köşe süsü
            sh._element.getparent().remove(sh._element)
    contained_accents(s)

# 2) Slayt 4 (Öne Çıkan Özellikler) + Slayt 5 (Mimari) kutu alt-yazılarını büyüt/koyulaştır
for idx in (3, 4):
    s = prs.slides[idx]
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None and r.font.size <= Pt(11.6):
                    r.font.size = Pt(13)
                    r.font.color.rgb = INK

prs.save(PATH)

# 3) Doğrulama: hiçbir döndürülmüş şekil slayt sınırını taşmasın
def overflow(sh):
    rot = abs(getattr(sh, "rotation", 0) or 0)
    if rot <= 0.5:
        return False
    a = math.radians(rot)
    hw = (abs(sh.width*math.cos(a)) + abs(sh.height*math.sin(a))) / 2
    hh = (abs(sh.width*math.sin(a)) + abs(sh.height*math.cos(a))) / 2
    cx, cy = sh.left + sh.width/2, sh.top + sh.height/2
    tol = 9000  # ~0.01 in tolerans
    return cx-hw < -tol or cy-hh < -tol or cx+hw > SW+tol or cy+hh > SH+tol

p2 = Presentation(PATH)
bad = 0
for i, s in enumerate(p2.slides, 1):
    if i in (1, 10):   # kapak/kapanış tam-taşma tasarımı (kasıtlı)
        continue
    for sh in s.shapes:
        if overflow(sh):
            bad += 1
            print("TASMA slayt", i)
print("OK — icerik slaytlarinda tasma:", bad)
