# -*- coding: utf-8 -*-
"""Köşe süslerini temiz ÜÇGEN KAMALARLA değiştirir (köşeye tam oturur, taşma yok).
Kapak: sol-üst + sağ-alt; içerik slaytları: sağ-üst + sağ-alt. Slayt metinlerine
(özellikle slayt 7 Ekstra Teknolojiler) DOKUNULMAZ. Yerinde çalışır."""
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PATH = "docs/Finans_Portali_Sunum.pptx"
NAVY = RGBColor(0x1E, 0x3A, 0x8A)
BLUE = RGBColor(0x1D, 0x4E, 0xD8)

prs = Presentation(PATH)
SW, SH = prs.slide_width, prs.slide_height


def wedge(s, corner, w, h, color):
    """corner: 'TL','TR','BL','BR' — o köşeye tam oturan dik üçgen (kama)."""
    W, Hh = Inches(w), Inches(h)
    left = 0 if corner[1] == "L" else SW - W
    top  = 0 if corner[0] == "T" else SH - Hh
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, left, top, W, Hh)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    xfrm = shp._element.spPr.find(qn("a:xfrm"))
    # Varsayılan dik açı SOL-ALT'ta. Üst köşe için flipV, sağ köşe için flipH.
    if corner[0] == "T": xfrm.set("flipV", "1")
    if corner[1] == "R": xfrm.set("flipH", "1")
    return shp


def strip_rotated(s):
    for sh in list(s.shapes):
        if abs(getattr(sh, "rotation", 0) or 0) > 0.5:
            sh._element.getparent().remove(sh._element)


# KAPAK (slayt 1): eğik bantları kaldır → sol-üst + sağ-alt büyük kamalar
s0 = prs.slides[0]
strip_rotated(s0)
wedge(s0, "TL", 3.3, 2.3, NAVY); wedge(s0, "TL", 2.15, 1.5, BLUE)
wedge(s0, "BR", 3.3, 2.3, NAVY); wedge(s0, "BR", 2.15, 1.5, BLUE)

# İÇERİK slaytları (2–9): eski döndürülmüş süsleri kaldır → sağ-üst + sağ-alt kamalar
for idx in range(1, 9):
    s = prs.slides[idx]
    strip_rotated(s)
    wedge(s, "TR", 1.85, 1.25, NAVY); wedge(s, "TR", 1.2, 0.8, BLUE)
    wedge(s, "BR", 1.85, 1.25, NAVY); wedge(s, "BR", 1.2, 0.8, BLUE)

prs.save(PATH)
print("OK — kamalar eklendi, taşma yok")
